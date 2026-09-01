from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer

from main import run_all

RESULTS = ROOT / "results"


def _comparison_charts(summary: dict) -> None:
    christ = summary["christofides"]["distance_km"]
    ga = summary["genetic"]["distance_km"]

    fig, ax = plt.subplots(figsize=(7, 5))
    ax.bar(["Christofides", "Génétique"], [christ, ga])
    ax.set_ylabel("Distance (km)")
    ax.set_title("Comparaison des distances")
    for i, value in enumerate([christ, ga]):
        ax.text(i, value + 20, f"{value:.1f} km", ha="center")
    fig.tight_layout()
    fig.savefig(RESULTS / "comparaison_distances.png", dpi=180)
    plt.close(fig)

    rows = summary["benchmark"]
    labels = [str(row["configuration"]) for row in rows]
    means = [float(row["mean_km"]) for row in rows]
    stds = [float(row["std_km"]) for row in rows]
    fig, ax = plt.subplots(figsize=(8, 5))
    ax.bar(labels, means, yerr=stds, capsize=5)
    ax.set_ylabel("Distance moyenne (km)")
    ax.set_title("Benchmark des configurations génétiques")
    fig.tight_layout()
    fig.savefig(RESULTS / "benchmark_ga.png", dpi=180)
    plt.close(fig)


def _add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.9), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.7), Inches(1.18), Inches(11.9), Inches(0.5))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)


def _add_bullets(slide, bullets: list[str], top: float = 1.8) -> None:
    box = slide.shapes.add_textbox(Inches(0.9), Inches(top), Inches(11.2), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.space_after = Pt(10)


def _make_pptx(summary: dict) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = [
        ("Le marchand ambulant", ["TSP sur 20 villes françaises", "Christofides vs algorithme génétique"]),
        ("Modélisation", ["20 villes issues du CSV fourni", "Graphe complet pondéré", "Distance de Haversine", "TSP métrique"]),
        ("Christofides", ["Prim → arbre couvrant minimal", "Sommets de degré impair", "Matching parfait minimum", "Hierholzer + shortcut", f"Distance : {summary['christofides']['distance_km']:.2f} km"]),
        ("Algorithme génétique", ["Individu = permutation des villes", "Sélection par tournoi", "Ordered Crossover (OX)", "Mutation swap / inversion", "Élitisme"]),
        ("Configuration retenue", ["Population : 160", "Générations : 520", "Tournoi : 5", "Crossover : 95%", "Mutation : 22%", "Élite : 6"]),
        ("Résultats", [f"Christofides : {summary['christofides']['distance_km']:.2f} km", f"Génétique : {summary['genetic']['distance_km']:.2f} km", f"Gain : {abs(summary['difference_ga_minus_christofides_km']):.2f} km", f"Amélioration : {abs(summary['difference_percent']):.2f}%"]),
        ("Analyse comparative", ["Christofides : rapide, déterministe, garantie ≤ 1,5× optimum", "Génétique : meilleure distance observée, mais stochastique", "Le GA demande davantage de paramétrage et plusieurs runs"]),
        ("Cartes et convergence", ["Carte interactive Christofides", "Carte interactive génétique", "Courbe de convergence du GA", "Benchmark de trois configurations"]),
        ("Organisation", ["Données / modélisation", "Christofides", "Algorithme génétique", "Benchmark / visualisation", "Documentation / présentation", "Voir TRELLO_ORGANISATION.md"]),
        ("Conclusion", ["Recommandation : GA si la priorité est la distance", "Christofides si la priorité est la rapidité et la garantie théorique", "3157,13 km = meilleure solution observée, pas preuve d'optimalité"]),
    ]

    for idx, (title, bullets) in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        _add_title(slide, title, f"Le marchand ambulant — {idx + 1}/10")
        _add_bullets(slide, bullets)
        if idx == 5 and (RESULTS / "comparaison_distances.png").exists():
            slide.shapes.add_picture(str(RESULTS / "comparaison_distances.png"), Inches(7.2), Inches(2.0), width=Inches(5.5))
        if idx == 7 and (RESULTS / "convergence_genetique.png").exists():
            slide.shapes.add_picture(str(RESULTS / "convergence_genetique.png"), Inches(7.2), Inches(2.0), width=Inches(5.5))

    prs.save(ROOT / "presentation_marchand_ambulant.pptx")


def _make_pdf(summary: dict) -> None:
    path = ROOT / "presentation_marchand_ambulant.pdf"
    doc = SimpleDocTemplate(
        str(path),
        pagesize=landscape(A4),
        rightMargin=1.5 * cm,
        leftMargin=1.5 * cm,
        topMargin=1.2 * cm,
        bottomMargin=1.2 * cm,
    )
    styles = getSampleStyleSheet()
    story = []
    content = [
        ("Le marchand ambulant — TSP", ["Projet sur 20 villes françaises", "Comparaison Christofides / algorithme génétique"]),
        ("Modélisation", ["Graphe complet", "Distances Haversine", "TSP métrique"]),
        ("Christofides", [f"Distance : {summary['christofides']['distance_km']:.2f} km", "Prim, matching minimum, Hierholzer, shortcut"]),
        ("Algorithme génétique", [f"Distance : {summary['genetic']['distance_km']:.2f} km", "Tournoi, OX, mutation, élitisme"]),
        ("Paramètres", ["Population 160", "520 générations", "Mutation 22%", "Crossover 95%", "Élite 6"]),
        ("Comparaison", [f"Gain GA : {abs(summary['difference_ga_minus_christofides_km']):.2f} km", f"Amélioration : {abs(summary['difference_percent']):.2f}%"]),
        ("Conclusion", ["GA recommandé pour la distance", "Christofides recommandé pour la rapidité et la garantie théorique", "Le GA ne prouve pas l'optimum global"]),
    ]
    for i, (title, bullets) in enumerate(content):
        story.append(Paragraph(title, styles["Title"]))
        story.append(Spacer(1, 0.5 * cm))
        for bullet in bullets:
            story.append(Paragraph(f"- {bullet}", styles["Heading2"]))
            story.append(Spacer(1, 0.25 * cm))
        if i != len(content) - 1:
            story.append(PageBreak())
    doc.build(story)


def main() -> None:
    summary = run_all()
    _comparison_charts(summary)
    _make_pptx(summary)
    _make_pdf(summary)
    print("Artifacts generated successfully")


if __name__ == "__main__":
    main()
